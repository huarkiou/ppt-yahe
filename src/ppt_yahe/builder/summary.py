"""Summary-slide builder with layout-config-driven parameter grouping."""

from __future__ import annotations

import logging

from pptx.chart.chart import Chart
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.presentation import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.slide import Slide
from pptx.table import Table
from pptx.util import Inches, Pt

from ppt_yahe.builder.types import SummaryLayoutConfig
from ppt_yahe.table_utils import create_styled_table, set_cell_text, set_merged_cell

_logger = logging.getLogger(__name__)

BLANK_SLIDE_LAYOUT = 6


def build_summary_slide(
    prs: Presentation,
    title: str,
    displacement_levels: list[str],
    section_ids: list[str],
    measurement_data: dict[tuple[str, str], tuple[float, float]],
    layout: SummaryLayoutConfig | None = None,
) -> None:
    """Build a summary slide with measurement table and comparison chart.

    All visual parameters are driven by *layout*; passing ``None`` uses
    the :class:`SummaryLayoutConfig` defaults.

    Args:
        prs: Presentation to add the slide to.
        title: Slide title.
        displacement_levels: Row labels (displacement levels).
        section_ids: Column labels (experiment sections).
        measurement_data: ``(displacement, section_id) → (force, length)``.
        layout: Visual and layout configuration.
    """
    cfg = layout or SummaryLayoutConfig()
    measurement_data_str = {k: (f"{v[0]:.2f}", f"{v[1]:.2f}") for k, v in measurement_data.items()}
    n_rows = len(displacement_levels)
    n_cols = len(section_ids)
    total_cols = 3 + n_cols
    total_rows = 2 + n_rows * 2 + 2

    slide: Slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height
    if slide_width_emu is None:
        raise ValueError("Slide width is required but got None")
    if slide_height_emu is None:
        raise ValueError("Slide height is required but got None")

    fixed = cfg.header_col_width + cfg.param_col_width + cfg.criteria_col_width
    total_table_inches = slide_width_emu / Inches(1)
    data_col_width = max(0, (total_table_inches - fixed) / n_cols) if n_cols else 0

    table_top = Inches(cfg.table_top_offset)
    table_height = Inches(
        cfg.data_row_height + cfg.header_row_height + (total_rows - 2) * cfg.data_row_height
    )

    table: Table = create_styled_table(
        slide,
        total_rows,
        total_cols,
        Inches(0),
        table_top,
        slide_width_emu,
        table_height,
    )

    table.columns[0].width = Inches(cfg.header_col_width)
    table.columns[1].width = Inches(cfg.param_col_width)
    table.columns[2].width = Inches(cfg.criteria_col_width)
    for c in range(3, total_cols):
        table.columns[c].width = Inches(data_col_width)

    table.rows[0].height = Inches(cfg.data_row_height)
    table.rows[1].height = Inches(cfg.header_row_height)
    for r in range(2, total_rows):
        table.rows[r].height = Inches(cfg.data_row_height)

    _populate_summary_table(
        table=table,
        title=title,
        displacement_levels=displacement_levels,
        section_ids=section_ids,
        measurement_data_str=measurement_data_str,
        total_cols=total_cols,
        n_rows=n_rows,
        cfg=cfg,
    )

    _add_comparison_chart(
        slide=slide,
        title=title,
        slide_width_emu=slide_width_emu,
        slide_height_emu=slide_height_emu,
        table_top=table_top,
        table_height=table_height,
        total_table_inches=total_table_inches,
        displacement_levels=displacement_levels,
        section_ids=section_ids,
        measurement_data=measurement_data,
        cfg=cfg,
    )


def _populate_summary_table(
    *,
    table: Table,
    title: str,
    displacement_levels: list[str],
    section_ids: list[str],
    measurement_data_str: dict[tuple[str, str], tuple[str, str]],
    total_cols: int,
    n_rows: int,
    cfg: SummaryLayoutConfig,
) -> None:
    """Fill all cells in the summary table."""
    # --- title / header rows ---
    # Left 3 non-data columns: merge row 0 vertically with row 1
    set_merged_cell(
        table, 0, 0, 1, 0, "类别", bold=True, font_size=cfg.font_size_header
    )
    set_merged_cell(
        table, 0, 1, 1, 1, "参数", bold=True, font_size=cfg.font_size_header
    )
    set_merged_cell(
        table, 0, 2, 1, 2, "评判标准", bold=True, font_size=cfg.font_size_header
    )

    # Data columns in row 0: merge horizontally for the title
    if total_cols > 3:
        set_merged_cell(
            table, 0, 3, 0, total_cols - 1, title,
            bold=True, font_size=cfg.font_size_title,
        )

    # Row 1 data column headers
    for j, section_id in enumerate(section_ids):
        set_cell_text(table.cell(1, 3 + j), section_id, bold=True, font_size=cfg.font_size_header)

    # --- force section ---
    force_start, force_end = 2, 1 + n_rows
    if n_rows > 1:
        set_merged_cell(
            table, force_start, 0, force_end, 0, "力值", bold=True, font_size=cfg.font_size_header
        )
    else:
        set_cell_text(table.cell(force_start, 0), "力值", bold=True, font_size=cfg.font_size_header)

    for i, displacement in enumerate(displacement_levels):
        row_idx = force_start + i
        set_cell_text(
            table.cell(row_idx, 1), displacement, bold=True, font_size=cfg.font_size_header
        )
        set_cell_text(table.cell(row_idx, 2), "", font_size=cfg.font_size_header)
        for j, section_id in enumerate(section_ids):
            info = measurement_data_str.get((displacement, section_id), ("", ""))
            set_cell_text(table.cell(row_idx, 3 + j), info[0], font_size=cfg.font_size_header)

    # --- length section ---
    length_start, length_end = 2 + n_rows, 1 + n_rows * 2
    if n_rows > 1:
        set_merged_cell(
            table, length_start, 0, length_end, 0, "长度", bold=True, font_size=cfg.font_size_header
        )
    else:
        set_cell_text(
            table.cell(length_start, 0), "长度", bold=True, font_size=cfg.font_size_header
        )

    for i, displacement in enumerate(displacement_levels):
        row_idx = length_start + i
        set_cell_text(
            table.cell(row_idx, 1), displacement, bold=True, font_size=cfg.font_size_header
        )
        set_cell_text(table.cell(row_idx, 2), "", font_size=cfg.font_size_header)
        for j, section_id in enumerate(section_ids):
            info = measurement_data_str.get((displacement, section_id), ("", ""))
            set_cell_text(table.cell(row_idx, 3 + j), info[1], font_size=cfg.font_size_header)

    # --- conclusion / remark rows ---
    conclusion_row = 2 + n_rows * 2
    set_cell_text(table.cell(conclusion_row, 0), "结论", bold=True, font_size=cfg.font_size_header)
    set_merged_cell(
        table, conclusion_row, 1, conclusion_row, total_cols - 1, "", font_size=cfg.font_size_header
    )

    remark_row = conclusion_row + 1
    set_cell_text(table.cell(remark_row, 0), "备注", bold=True, font_size=cfg.font_size_header)
    set_merged_cell(
        table, remark_row, 1, remark_row, total_cols - 1, "", font_size=cfg.font_size_header
    )


def _add_comparison_chart(
    *,
    slide: Slide,
    title: str,
    slide_width_emu: int,
    slide_height_emu: int,
    table_top: int,
    table_height: int,
    total_table_inches: float,
    displacement_levels: list[str],
    section_ids: list[str],
    measurement_data: dict[tuple[str, str], tuple[float, float]],
    cfg: SummaryLayoutConfig,
) -> None:
    """Add a clustered column chart below the summary table."""
    table_bottom = (table_top + table_height) / Inches(1)
    chart_top_inch = table_bottom + cfg.chart_top_padding
    slide_height_inch = slide_height_emu / Inches(1)
    max_chart_height = slide_height_inch - chart_top_inch - cfg.chart_bottom_padding
    chart_height_inch = min(max_chart_height, cfg.chart_max_height)
    chart_left_inch = cfg.chart_left_indent
    chart_width_inch = total_table_inches - 2 * chart_left_inch

    chart_data = CategoryChartData()
    chart_data.categories = displacement_levels
    for section_id in section_ids:
        values = [measurement_data.get((d, section_id), (0, 0))[0] for d in displacement_levels]
        chart_data.add_series(section_id, values)

    chart_frame: GraphicFrame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(chart_left_inch),
        Inches(chart_top_inch),
        Inches(chart_width_inch),
        Inches(chart_height_inch),
        chart_data,  # ty:ignore[invalid-argument-type]
    )  # ty:ignore[invalid-assignment]

    chart: Chart = chart_frame.chart

    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(cfg.font_size_chart)

    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(cfg.font_size_chart)

    chart.category_axis.tick_labels.font.size = Pt(cfg.font_size_chart)

    chart.value_axis.has_major_gridlines = False
    chart.value_axis.tick_labels.font.size = Pt(cfg.font_size_chart)

    plot = chart.plots[0]
    plot.has_data_labels = True
    for series in chart.series:
        series_data_labels = series.data_labels
        series_data_labels.show_value = True
        series_data_labels.number_format = "0.00"
        series_data_labels.font.size = Pt(cfg.font_size_chart)
        series_data_labels.font.name = "微软雅黑"
