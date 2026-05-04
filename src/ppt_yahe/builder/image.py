"""Image-matrix slide builder with layout-config-driven parameter grouping."""

from __future__ import annotations

import logging
from pathlib import Path

from PIL import Image
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.table import Table
from pptx.util import Inches, Pt

from ppt_yahe.builder.types import CellDimensions, ImageLayoutConfig
from ppt_yahe.table_utils import (
    column_left_inch,
    create_styled_table,
    get_measurement_str,
    row_top_inch,
    set_cell_text,
    set_merged_cell,
)

_logger = logging.getLogger(__name__)

BLANK_SLIDE_LAYOUT = 6


def build_image_slide(
    prs: Presentation,
    image_dir: str | Path,
    displacement_levels: list[str],
    section_ids: list[str],
    filename_template: str = "{displacement}_{section}",
    top_left_label: str = "",
    measurement_data: dict[tuple[str, str], tuple[float, float]] | None = None,
    layout: ImageLayoutConfig | None = None,
    image_fill_uniform: bool = False,
) -> None:
    """Build an image-matrix slide with experiment images.

    All visual parameters are driven by *layout*; passing ``None`` uses
    the :class:`ImageLayoutConfig` defaults.

    Args:
        prs: Presentation to add the slide to.
        image_dir: Directory containing experiment images.
        displacement_levels: Row labels (displacement levels).
        section_ids: Column labels (experiment sections).
        filename_template: Template for image filenames (supports
            ``{displacement}`` and ``{section}``).
        top_left_label: Label for the top-left corner cell.
        measurement_data: ``(displacement, section_id) → (force, length)``.
        layout: Visual and layout configuration.
        image_fill_uniform: Fit images uniformly if ``True``; fill cell otherwise.
    """
    cfg = layout or ImageLayoutConfig()
    image_dir = Path(image_dir)
    if measurement_data is None:
        measurement_data = {}

    slide, usable_width, usable_height, n_rows, n_cols = _validate_and_init_layout(
        prs, displacement_levels, section_ids, cfg
    )

    cell_dim = _compute_cell_dimensions(n_rows, n_cols, usable_width, usable_height, cfg)

    table, table_left, table_top = _setup_table_grid(
        slide=slide,
        usable_width=usable_width,
        usable_height=usable_height,
        n_rows=n_rows,
        cell_dim=cell_dim,
        cfg=cfg,
    )

    _populate_headers(table, displacement_levels, section_ids, top_left_label, cfg)
    _populate_measurement_cells(table, displacement_levels, section_ids, measurement_data, cfg)

    _insert_images(
        slide=slide,
        table=table,
        displacement_levels=displacement_levels,
        section_ids=section_ids,
        image_dir=image_dir,
        filename_template=filename_template,
        image_fill_uniform=image_fill_uniform,
        table_left=table_left,
        table_top=table_top,
        cell_dim=cell_dim,
    )


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------


def _validate_and_init_layout(
    prs: Presentation,
    displacement_levels: list[str],
    section_ids: list[str],
    cfg: ImageLayoutConfig,
) -> tuple[Slide, float, float, int, int]:
    """Validate parameters, create the slide, and compute usable area."""
    if prs.slide_width is None or prs.slide_height is None:
        raise ValueError("Slide dimensions are required but got None")

    slide_width: float = prs.slide_width / Inches(1)
    slide_height: float = prs.slide_height / Inches(1)

    usable_width = slide_width - 2 * cfg.margin_lr
    usable_height = slide_height - 2 * cfg.margin_tb

    n_rows = len(displacement_levels)
    n_cols = len(section_ids)
    if n_rows == 0 or n_cols == 0:
        raise ValueError("displacement_levels and section_ids must not be empty")

    slide: Slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
    return (slide, usable_width, usable_height, n_rows, n_cols)


def _compute_cell_dimensions(
    n_rows: int,
    n_cols: int,
    usable_width: float,
    usable_height: float,
    cfg: ImageLayoutConfig,
) -> CellDimensions:
    """Compute all cell dimension values from layout parameters.

    Returns a typed :class:`CellDimensions` object instead of an
    opaque dict.
    """
    data_area_width = usable_width - cfg.header_col_width
    data_area_height = usable_height - cfg.header_row_height

    total_data_cols = n_cols * 2
    total_data_rows = n_rows * 3

    sub_cell_width = data_area_width / total_data_cols
    sub_cell_height = data_area_height / total_data_rows

    three_row_height = 3 * sub_cell_height
    supplement_total_height = three_row_height * cfg.supplement_row_ratio
    image_row_height = three_row_height - supplement_total_height
    title_row_height = supplement_total_height / 2
    data_row_height = supplement_total_height / 2

    image_area_width = sub_cell_width * 2
    image_area_height = image_row_height

    padding = Pt(cfg.image_padding_pt).inches
    available_width = image_area_width - 2 * padding
    available_height = image_area_height - 2 * padding
    square_size = min(available_width, available_height)

    return CellDimensions(
        data_area_width=data_area_width,
        data_area_height=data_area_height,
        total_data_cols=total_data_cols,
        total_data_rows=total_data_rows,
        total_cols=1 + total_data_cols,
        total_rows=1 + total_data_rows,
        sub_cell_width=sub_cell_width,
        sub_cell_height=sub_cell_height,
        three_row_height=three_row_height,
        supplement_total_height=supplement_total_height,
        image_row_height=image_row_height,
        title_row_height=title_row_height,
        data_row_height=data_row_height,
        image_area_width=image_area_width,
        image_area_height=image_area_height,
        padding=padding,
        available_width=available_width,
        available_height=available_height,
        square_size=square_size,
    )


def _setup_table_grid(
    *,
    slide: Slide,
    usable_width: float,
    usable_height: float,
    n_rows: int,
    cell_dim: CellDimensions,
    cfg: ImageLayoutConfig,
) -> tuple[Table, float, float]:
    """Create and configure the table grid on the slide."""
    table = create_styled_table(
        slide,
        cell_dim.total_rows,
        cell_dim.total_cols,
        Inches(cfg.margin_lr),
        Inches(cfg.margin_tb),
        Inches(usable_width),
        Inches(usable_height),
    )

    table_left = cfg.margin_lr
    table_top = cfg.margin_tb

    table.columns[0].width = Inches(cfg.header_col_width)
    for c in range(1, cell_dim.total_cols):
        table.columns[c].width = Inches(cell_dim.sub_cell_width)

    table.rows[0].height = Inches(cfg.header_row_height)
    for i in range(n_rows):
        base = 1 + i * 3
        table.rows[base].height = Inches(cell_dim.image_row_height)
        table.rows[base + 1].height = Inches(cell_dim.title_row_height)
        table.rows[base + 2].height = Inches(cell_dim.data_row_height)

    return (table, table_left, table_top)


def _populate_headers(
    table: Table,
    displacement_levels: list[str],
    section_ids: list[str],
    top_left_label: str,
    cfg: ImageLayoutConfig,
) -> None:
    """Fill the header row and first column with labels."""
    set_cell_text(
        table.cell(0, 0),
        top_left_label,
        bold=True,
        font_size=cfg.font_size_header,
    )

    for j, section_id in enumerate(section_ids):
        col_a = 1 + j * 2
        col_b = col_a + 1
        set_merged_cell(
            table,
            0,
            col_a,
            0,
            col_b,
            section_id,
            bold=True,
            font_size=cfg.font_size_header,
        )

    for i, displacement in enumerate(displacement_levels):
        row_a = 1 + i * 3
        row_b = row_a + 2
        set_merged_cell(
            table,
            row_a,
            0,
            row_b,
            0,
            displacement,
            bold=True,
            font_size=cfg.font_size_header,
        )


def _populate_measurement_cells(
    table: Table,
    displacement_levels: list[str],
    section_ids: list[str],
    measurement_data: dict[tuple[str, str], tuple[float, float]],
    cfg: ImageLayoutConfig,
) -> None:
    """Fill the measurement data cells (title and data rows)."""
    for i, displacement in enumerate(displacement_levels):
        for j, section_id in enumerate(section_ids):
            image_row = 1 + i * 3
            title_row = image_row + 1
            data_row = image_row + 2
            col_a = 1 + j * 2
            col_b = col_a + 1

            # merge image cell across two sub-columns
            table.cell(image_row, col_a).merge(table.cell(image_row, col_b))

            set_cell_text(
                table.cell(title_row, col_a),
                "力值",
                bold=True,
                font_size=cfg.font_size_data,
            )
            set_cell_text(
                table.cell(title_row, col_b),
                "长度",
                bold=True,
                font_size=cfg.font_size_data,
            )

            left_text, right_text = get_measurement_str(
                measurement_data,
                displacement,
                section_id,
            )
            set_cell_text(
                table.cell(data_row, col_a),
                left_text,
                font_size=cfg.font_size_data,
            )
            set_cell_text(
                table.cell(data_row, col_b),
                right_text,
                font_size=cfg.font_size_data,
            )


def _insert_images(
    *,
    slide: Slide,
    table: Table,
    displacement_levels: list[str],
    section_ids: list[str],
    image_dir: Path,
    filename_template: str,
    image_fill_uniform: bool,
    table_left: float,
    table_top: float,
    cell_dim: CellDimensions,
) -> None:
    """Insert images into the table cells using computed dimensions."""
    square_size = cell_dim.square_size
    available_width = cell_dim.available_width
    available_height = cell_dim.available_height

    for i, displacement in enumerate(displacement_levels):
        for j, section_id in enumerate(section_ids):
            filename = filename_template.format(
                displacement=displacement,
                section=section_id,
            )
            img_path = image_dir / filename
            for suffix in (".png", ".jpg"):
                candidate = img_path.with_suffix(suffix)
                if candidate.exists():
                    img_path = candidate
                    break
            else:
                _logger.warning(
                    "跳过不存在的图片: %s.[png|jpg]",
                    img_path,
                )
                continue

            try:
                with Image.open(img_path) as im:
                    original_width, original_height = im.size
            except Exception:
                _logger.error("无法打开图片: %s", img_path)
                continue

            aspect = original_width / original_height

            if image_fill_uniform:
                if aspect >= 1:
                    display_width = square_size
                    display_height = square_size / aspect
                else:
                    display_height = square_size
                    display_width = square_size * aspect
            else:
                if available_width / available_height >= aspect:
                    display_height = available_height
                    display_width = available_height * aspect
                else:
                    display_width = available_width
                    display_height = available_width / aspect

            image_row = 1 + i * 3
            col_a = 1 + j * 2

            cell_left = table_left + column_left_inch(table, col_a)
            cell_top = table_top + row_top_inch(table, image_row)
            cell_w = (table.columns[col_a].width + table.columns[col_a + 1].width) / Inches(1)
            cell_h = table.rows[image_row].height / Inches(1)

            offset_x = (cell_w - display_width) / 2
            offset_y = (cell_h - display_height) / 2

            left = Inches(cell_left + offset_x)
            top = Inches(cell_top + offset_y)

            try:
                slide.shapes.add_picture(
                    str(img_path),
                    left=left,
                    top=top,
                    width=Inches(display_width),
                )
            except Exception:
                _logger.error("无法插入图片: %s", img_path)
