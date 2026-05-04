import logging

from pptx.chart.chart import Chart
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.presentation import Presentation
from pptx.shapes.graphfrm import GraphicFrame
from pptx.slide import Slide
from pptx.util import Inches, Pt

from ppt_yahe.table_utils import create_styled_table, set_cell_text, set_merged_cell

logger = logging.getLogger(__name__)

# Named constants for layout dimensions
DEFAULT_HEADER_COL_WIDTH = 0.9
DEFAULT_PARAM_COL_WIDTH = 1.0
DEFAULT_CRITERIA_COL_WIDTH = 1.2
DEFAULT_HEADER_ROW_HEIGHT = 0.4
DEFAULT_DATA_ROW_HEIGHT = 0.3
FONT_SIZE_TITLE = 14
FONT_SIZE_HEADER = 12
FONT_SIZE_CHART = 10
CHART_TOP_PADDING = 0.0
CHART_BOTTOM_PADDING = 0.2
CHART_MAX_HEIGHT = 3.5
CHART_LEFT_INDENT = 1.6
TABLE_TOP_OFFSET = 0.4
BLANK_SLIDE_LAYOUT = 6


def add_summary_slide(
    prs: Presentation,
    title: str,
    displacement_levels: list[str],
    section_ids: list[str],
    measurement_data: dict[tuple[str, str], tuple[float, float]],
    header_col_width: float = DEFAULT_HEADER_COL_WIDTH,
    param_col_width: float = DEFAULT_PARAM_COL_WIDTH,
    criteria_col_width: float = DEFAULT_CRITERIA_COL_WIDTH,
    header_row_height: float = DEFAULT_HEADER_ROW_HEIGHT,
    data_row_height: float = DEFAULT_DATA_ROW_HEIGHT,
) -> None:
    """Add a summary slide with measurement data table and comparison chart.

    Creates a slide containing a structured table of measurement data organized
    by displacement levels and section IDs, followed by a bar chart comparing
    force values across sections.

    Args:
        prs: The PowerPoint presentation object to add the slide to.
        title: Title displayed in the table header row and chart title.
        displacement_levels: List of displacement level labels for data rows.
        section_ids: List of section identifiers used as column headers.
        measurement_data: Dictionary mapping (displacement_level, section_id) tuples
            to (force_value, length_value) measurement pairs.
        header_col_width: Width of the header column in inches (default: 0.9).
        param_col_width: Width of the parameter column in inches (default: 1.0).
        criteria_col_width: Width of the criteria column in inches (default: 1.2).
        header_row_height: Height of the header row in inches (default: 0.4).
        data_row_height: Height of data rows in inches (default: 0.3).
    """
    measurement_data_str = {k: (str(v[0]), str(v[1])) for k, v in measurement_data.items()}
    n_rows = len(displacement_levels)
    n_cols = len(section_ids)
    total_cols = 3 + n_cols
    total_rows = 2 + n_rows * 2 + 2

    slide: Slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height
    if slide_width_emu is None:
        raise ValueError("Slide width is required but got None")
    if slide_height_emu is None:
        raise ValueError("Slide height is required but got None")

    fixed = header_col_width + param_col_width + criteria_col_width
    total_table_inches = slide_width_emu / Inches(1)
    data_col_width = max(0, (total_table_inches - fixed) / n_cols) if n_cols else 0

    table_top = Inches(TABLE_TOP_OFFSET)
    table_height = Inches(data_row_height + header_row_height + (total_rows - 2) * data_row_height)

    table = create_styled_table(
        slide,
        total_rows,
        total_cols,
        Inches(0),
        table_top,
        slide_width_emu,
        table_height,
    )

    table.columns[0].width = Inches(header_col_width)
    table.columns[1].width = Inches(param_col_width)
    table.columns[2].width = Inches(criteria_col_width)
    for c in range(3, total_cols):
        table.columns[c].width = Inches(data_col_width)

    table.rows[0].height = Inches(data_row_height)
    table.rows[1].height = Inches(header_row_height)
    for r in range(2, total_rows):
        table.rows[r].height = Inches(data_row_height)

    set_merged_cell(table, 0, 0, 0, total_cols - 1, title, bold=True, font_size=FONT_SIZE_TITLE)
    set_cell_text(table.cell(1, 0), "类别", bold=True, font_size=FONT_SIZE_HEADER)
    set_cell_text(table.cell(1, 1), "参数", bold=True, font_size=FONT_SIZE_HEADER)
    set_cell_text(table.cell(1, 2), "评判标准", bold=True, font_size=FONT_SIZE_HEADER)
    for j, section_id in enumerate(section_ids):
        set_cell_text(table.cell(1, 3 + j), section_id, bold=True, font_size=FONT_SIZE_HEADER)

    force_start, force_end = 2, 1 + n_rows
    if n_rows > 1:
        set_merged_cell(
            table, force_start, 0, force_end, 0, "力值", bold=True, font_size=FONT_SIZE_HEADER
        )
    else:
        set_cell_text(table.cell(force_start, 0), "力值", bold=True, font_size=FONT_SIZE_HEADER)

    for i, displacement in enumerate(displacement_levels):
        row_idx = force_start + i
        set_cell_text(table.cell(row_idx, 1), displacement, bold=True, font_size=FONT_SIZE_HEADER)
        set_cell_text(table.cell(row_idx, 2), "", font_size=FONT_SIZE_HEADER)
        for j, section_id in enumerate(section_ids):
            info = measurement_data_str.get((displacement, section_id), ("", ""))
            set_cell_text(table.cell(row_idx, 3 + j), info[0], font_size=FONT_SIZE_HEADER)

    length_start, length_end = 2 + n_rows, 1 + n_rows * 2
    if n_rows > 1:
        set_merged_cell(
            table, length_start, 0, length_end, 0, "长度", bold=True, font_size=FONT_SIZE_HEADER
        )
    else:
        set_cell_text(table.cell(length_start, 0), "长度", bold=True, font_size=FONT_SIZE_HEADER)

    for i, displacement in enumerate(displacement_levels):
        row_idx = length_start + i
        set_cell_text(table.cell(row_idx, 1), displacement, bold=True, font_size=FONT_SIZE_HEADER)
        set_cell_text(table.cell(row_idx, 2), "", font_size=FONT_SIZE_HEADER)
        for j, section_id in enumerate(section_ids):
            info = measurement_data_str.get((displacement, section_id), ("", ""))
            set_cell_text(table.cell(row_idx, 3 + j), info[1], font_size=FONT_SIZE_HEADER)

    conclusion_row = 2 + n_rows * 2
    set_cell_text(table.cell(conclusion_row, 0), "结论", bold=True, font_size=FONT_SIZE_HEADER)
    set_merged_cell(
        table, conclusion_row, 1, conclusion_row, total_cols - 1, "", font_size=FONT_SIZE_HEADER
    )

    remark_row = conclusion_row + 1
    set_cell_text(table.cell(remark_row, 0), "备注", bold=True, font_size=FONT_SIZE_HEADER)
    set_merged_cell(
        table, remark_row, 1, remark_row, total_cols - 1, "", font_size=FONT_SIZE_HEADER
    )

    _add_comparison_chart(
        slide,
        title,
        slide_width_emu,
        slide_height_emu,
        table_top,
        table_height,
        total_table_inches,
        displacement_levels,
        section_ids,
        measurement_data,
    )


def _add_comparison_chart(
    slide: Slide,
    title: str,
    slide_width_emu: int,
    slide_height_emu: int,
    table_top: int,
    table_height: int,
    total_table_inches: float,
    displacement_levels: list[str],
    section_ids: list[str],
    measurement_data: dict[tuple[str, str], tuple[float, float]],
) -> None:
    """Add a clustered column chart comparing force values across sections.

    Creates a bar chart positioned below the summary table, showing force
    measurements for each displacement level grouped by section.

    Args:
        slide: The slide to add the chart to.
        title: Chart title text.
        slide_width_emu: Slide width in EMU (English Metric Units).
        slide_height_emu: Slide height in EMU.
        table_top: Top position of the summary table in EMU.
        table_height: Height of the summary table in EMU.
        total_table_inches: Total width of the table in inches.
        displacement_levels: List of displacement level labels (chart categories).
        section_ids: List of section identifiers (chart series).
        measurement_data: Dictionary mapping (displacement_level, section_id) tuples
            to (force_value, length_value) measurement pairs. Only force values
            are plotted in the chart.
    """
    table_bottom = (table_top + table_height) / Inches(1)
    chart_top_inch = table_bottom + CHART_TOP_PADDING
    slide_height_inch = slide_height_emu / Inches(1)
    max_chart_height = slide_height_inch - chart_top_inch - CHART_BOTTOM_PADDING
    chart_height_inch = min(max_chart_height, CHART_MAX_HEIGHT)
    chart_left_inch = CHART_LEFT_INDENT
    chart_width_inch = total_table_inches - 2 * chart_left_inch

    chart_data = CategoryChartData()
    chart_data.categories = displacement_levels
    for section_id in section_ids:
        values = [measurement_data.get((d, section_id), (0, 0))[0] for d in displacement_levels]
        chart_data.add_series(section_id, values)

    chart_frame: GraphicFrame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(chart_left_inch),
        Inches(chart_top_inch),
        Inches(chart_width_inch),
        Inches(chart_height_inch),
        chart_data,  # ty:ignore[invalid-argument-type]
    )  # ty:ignore[invalid-assignment]

    chart: Chart = chart_frame.chart

    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(FONT_SIZE_CHART)

    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(FONT_SIZE_CHART)

    chart.category_axis.tick_labels.font.size = Pt(FONT_SIZE_CHART)

    chart.value_axis.has_major_gridlines = False
    chart.value_axis.tick_labels.font.size = Pt(FONT_SIZE_CHART)

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_value = True
    data_labels.font.size = Pt(FONT_SIZE_CHART)
    for series in chart.series:
        for point in series.points:
            point.data_label.font.size = Pt(FONT_SIZE_CHART)
            point.data_label.font.name = "微软雅黑"
