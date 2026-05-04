import logging
from pathlib import Path

from PIL import Image
from pptx.presentation import Presentation
from pptx.slide import Slide
from pptx.table import Table
from pptx.util import Inches, Pt

from ppt_yahe.table_utils import (
    column_left_inch,
    create_styled_table,
    get_measurement_str,
    row_top_inch,
    set_cell_text,
    set_merged_cell,
)

_logger = logging.getLogger(__name__)

# Named constants for magic numbers
DEFAULT_MARGIN_LR = 0.0
DEFAULT_MARGIN_TB = 0.5
DEFAULT_HEADER_COL_WIDTH = 1.2
DEFAULT_HEADER_ROW_HEIGHT = 0.5
DEFAULT_FONT_SIZE_HEADER = 12
DEFAULT_FONT_SIZE_DATA = 10
DEFAULT_SUPPLEMENT_ROW_RATIO = 0.30
IMAGE_PADDING_PT = 1
BLANK_SLIDE_LAYOUT = 6


def add_image_slide(
    prs: Presentation,
    image_dir: str | Path,
    displacement_levels: list[str],
    section_ids: list[str],
    filename_template: str = "{displacement}_{section}",
    top_left_label: str = "",
    header_col_width: float = DEFAULT_HEADER_COL_WIDTH,
    header_row_height: float = DEFAULT_HEADER_ROW_HEIGHT,
    measurement_data: dict[tuple[str, str], tuple[float, float]] | None = None,
    supplement_row_ratio: float = DEFAULT_SUPPLEMENT_ROW_RATIO,
    image_fill_uniform: bool = False,
) -> None:
    """Add an image matrix slide to the presentation.

    Creates a slide with a grid layout containing experiment images
    arranged by displacement levels and section IDs, with measurement
    labels for force and length values.

    Args:
        prs: Presentation object to add the slide to.
        image_dir: Directory containing experiment images.
        displacement_levels: List of displacement level labels.
        section_ids: List of section identifiers.
        filename_template: Template string for image filenames.
            Supports ``{displacement}`` and ``{section}`` placeholders.
        top_left_label: Label for the top-left corner cell.
        header_col_width: Width of the header column in inches.
        header_row_height: Height of the header row in inches.
        measurement_data: Dict mapping ``(displacement, section_id)`` tuples
            to ``(force, length)`` measurement pairs.
        supplement_row_ratio: Ratio of the supplement row height relative to
            the total three-row block height.
        image_fill_uniform: If ``True``, fit images uniformly maintaining
            aspect ratio; if ``False``, fill available space.
    """
    image_dir = Path(image_dir)
    if measurement_data is None:
        measurement_data = {}

    slide, usable_width, usable_height, margin_lr, margin_tb, n_rows, n_cols = (
        _validate_and_init_layout(
            prs,
            displacement_levels,
            section_ids,
            header_col_width,
            header_row_height,
            supplement_row_ratio,
        )
    )

    cell_dim = _compute_cell_dimensions(
        n_rows,
        n_cols,
        usable_width,
        usable_height,
        header_col_width,
        header_row_height,
        supplement_row_ratio,
    )

    table, table_left, table_top = _setup_table_grid(
        slide,
        cell_dim["total_rows"],  # ty:ignore[invalid-argument-type]
        cell_dim["total_cols"],  # ty:ignore[invalid-argument-type]
        margin_lr,
        margin_tb,
        usable_width,
        usable_height,
        header_col_width,
        header_row_height,
        cell_dim["sub_cell_width"],
        n_rows,
        cell_dim["image_row_height"],
        cell_dim["title_row_height"],
        cell_dim["data_row_height"],
    )

    _populate_headers(table, displacement_levels, section_ids, top_left_label)

    _populate_measurement_cells(
        table,
        displacement_levels,
        section_ids,
        measurement_data,
    )

    _insert_images(
        slide,
        table,
        displacement_levels,
        section_ids,
        image_dir,
        filename_template,
        image_fill_uniform,
        table_left,
        table_top,
        cell_dim,
    )


def _validate_and_init_layout(
    prs: Presentation,
    displacement_levels: list[str],
    section_ids: list[str],
    header_col_width: float,
    header_row_height: float,
    supplement_row_ratio: float,
) -> tuple[Slide, float, float, float, float, int, int]:
    """Validate parameters, create the slide, and compute usable area.

    Args:
        prs: Presentation object.
        displacement_levels: List of displacement level labels.
        section_ids: List of section identifiers.
        header_col_width: Width of the header column in inches.
        header_row_height: Height of the header row in inches.
        supplement_row_ratio: Ratio of supplement row to three-row block.

    Returns:
        Tuple of (slide, usable_width, usable_height, margin_lr, margin_tb,
        n_rows, n_cols).
    """
    if prs.slide_width is None or prs.slide_height is None:
        raise ValueError("Slide dimensions are required but got None")

    slide_width: float = prs.slide_width / Inches(1)
    slide_height: float = prs.slide_height / Inches(1)

    margin_lr = DEFAULT_MARGIN_LR
    margin_tb = DEFAULT_MARGIN_TB
    usable_width = slide_width - 2 * margin_lr
    usable_height = slide_height - 2 * margin_tb

    n_rows = len(displacement_levels)
    n_cols = len(section_ids)
    if n_rows == 0 or n_cols == 0:
        raise ValueError("参数列表不能为空")

    slide: Slide = prs.slides.add_slide(prs.slide_layouts[BLANK_SLIDE_LAYOUT])

    return (slide, usable_width, usable_height, margin_lr, margin_tb, n_rows, n_cols)


def _compute_cell_dimensions(
    n_rows: int,
    n_cols: int,
    usable_width: float,
    usable_height: float,
    header_col_width: float,
    header_row_height: float,
    supplement_row_ratio: float,
) -> dict[str, float | int]:
    """Compute all cell dimension values from layout parameters.

    Args:
        n_rows: Number of data rows.
        n_cols: Number of data columns.
        usable_width: Usable slide width in inches.
        usable_height: Usable slide height in inches.
        header_col_width: Width of the header column in inches.
        header_row_height: Height of the header row in inches.
        supplement_row_ratio: Ratio of supplement row to three-row block.

    Returns:
        Dict containing all computed dimension values (sub_cell_width,
        sub_cell_height, image_row_height, title_row_height, data_row_height,
        square_size, available_width, available_height, padding, total_rows,
        total_cols, and intermediate values).
    """
    data_area_width = usable_width - header_col_width
    data_area_height = usable_height - header_row_height

    total_data_cols = n_cols * 2
    total_data_rows = n_rows * 3

    sub_cell_width = data_area_width / total_data_cols
    sub_cell_height = data_area_height / total_data_rows

    three_row_height = 3 * sub_cell_height
    supplement_total_height = three_row_height * supplement_row_ratio
    image_row_height = three_row_height - supplement_total_height
    title_row_height = supplement_total_height / 2
    data_row_height = supplement_total_height / 2

    image_area_width = sub_cell_width * 2
    image_area_height = image_row_height

    padding = Pt(IMAGE_PADDING_PT).inches
    available_width = image_area_width - 2 * padding
    available_height = image_area_height - 2 * padding
    square_size = min(available_width, available_height)

    total_cols = 1 + total_data_cols
    total_rows = 1 + total_data_rows

    return {
        "data_area_width": data_area_width,
        "data_area_height": data_area_height,
        "total_data_cols": total_data_cols,
        "total_data_rows": total_data_rows,
        "sub_cell_width": sub_cell_width,
        "sub_cell_height": sub_cell_height,
        "three_row_height": three_row_height,
        "supplement_total_height": supplement_total_height,
        "image_row_height": image_row_height,
        "title_row_height": title_row_height,
        "data_row_height": data_row_height,
        "image_area_width": image_area_width,
        "image_area_height": image_area_height,
        "padding": padding,
        "available_width": available_width,
        "available_height": available_height,
        "square_size": square_size,
        "total_cols": total_cols,
        "total_rows": total_rows,
    }


def _setup_table_grid(
    slide: Slide,
    total_rows: int,
    total_cols: int,
    margin_lr: float,
    margin_tb: float,
    usable_width: float,
    usable_height: float,
    header_col_width: float,
    header_row_height: float,
    sub_cell_width: float,
    n_rows: int,
    image_row_height: float,
    title_row_height: float,
    data_row_height: float,
) -> tuple[Table, float, float]:
    """Create and configure the table grid on the slide.

    Args:
        slide: Slide to add the table to.
        total_rows: Total number of rows in the table.
        total_cols: Total number of columns in the table.
        margin_lr: Left/right margin in inches.
        margin_tb: Top/bottom margin in inches.
        usable_width: Usable slide width in inches.
        usable_height: Usable slide height in inches.
        header_col_width: Width of the header column in inches.
        header_row_height: Height of the header row in inches.
        sub_cell_width: Width of each sub-cell in inches.
        n_rows: Number of data rows (excluding header).
        image_row_height: Height of image rows in inches.
        title_row_height: Height of title rows in inches.
        data_row_height: Height of data rows in inches.

    Returns:
        Tuple of (table, table_left, table_top).
    """
    table = create_styled_table(
        slide,
        total_rows,
        total_cols,
        Inches(margin_lr),
        Inches(margin_tb),
        Inches(usable_width),
        Inches(usable_height),
    )

    table_left = margin_lr
    table_top = margin_tb

    table.columns[0].width = Inches(header_col_width)
    for c in range(1, total_cols):
        table.columns[c].width = Inches(sub_cell_width)

    table.rows[0].height = Inches(header_row_height)
    for i in range(n_rows):
        base = 1 + i * 3
        table.rows[base].height = Inches(image_row_height)
        table.rows[base + 1].height = Inches(title_row_height)
        table.rows[base + 2].height = Inches(data_row_height)

    return (table, table_left, table_top)


def _populate_headers(
    table: Table,
    displacement_levels: list[str],
    section_ids: list[str],
    top_left_label: str,
) -> None:
    """Fill the header row and first column with labels.

    Args:
        table: The table to populate.
        displacement_levels: List of displacement level labels.
        section_ids: List of section identifiers.
        top_left_label: Label for the top-left corner cell.
    """
    set_cell_text(
        table.cell(0, 0),
        top_left_label,
        bold=True,
        font_size=DEFAULT_FONT_SIZE_HEADER,
    )

    for j, section_id in enumerate(section_ids):
        col_a = 1 + j * 2
        col_b = col_a + 1
        set_merged_cell(
            table,
            0,
            col_a,
            0,
            col_b,
            section_id,
            bold=True,
            font_size=DEFAULT_FONT_SIZE_HEADER,
        )

    for i, displacement in enumerate(displacement_levels):
        row_a = 1 + i * 3
        row_b = row_a + 2
        set_merged_cell(
            table,
            row_a,
            0,
            row_b,
            0,
            displacement,
            bold=True,
            font_size=DEFAULT_FONT_SIZE_HEADER,
        )


def _populate_measurement_cells(
    table: Table,
    displacement_levels: list[str],
    section_ids: list[str],
    measurement_data: dict[tuple[str, str], tuple[float, float]],
) -> None:
    """Fill the measurement data cells (title and data rows).

    Args:
        table: The table to populate.
        displacement_levels: List of displacement level labels.
        section_ids: List of section identifiers.
        measurement_data: Dict mapping (displacement, section_id) tuples
            to (force, length) measurement pairs.
    """
    for i, displacement in enumerate(displacement_levels):
        for j, section_id in enumerate(section_ids):
            image_row = 1 + i * 3
            title_row = image_row + 1
            data_row = image_row + 2
            col_a = 1 + j * 2
            col_b = col_a + 1

            table.cell(image_row, col_a).merge(table.cell(image_row, col_b))

            set_cell_text(
                table.cell(title_row, col_a),
                "力值",
                bold=True,
                font_size=DEFAULT_FONT_SIZE_DATA,
            )
            set_cell_text(
                table.cell(title_row, col_b),
                "长度",
                bold=True,
                font_size=DEFAULT_FONT_SIZE_DATA,
            )

            left_text, right_text = get_measurement_str(
                measurement_data,
                displacement,
                section_id,
            )
            set_cell_text(
                table.cell(data_row, col_a),
                left_text,
                font_size=DEFAULT_FONT_SIZE_DATA,
            )
            set_cell_text(
                table.cell(data_row, col_b),
                right_text,
                font_size=DEFAULT_FONT_SIZE_DATA,
            )


def _insert_images(
    slide: Slide,
    table: Table,
    displacement_levels: list[str],
    section_ids: list[str],
    image_dir: Path,
    filename_template: str,
    image_fill_uniform: bool,
    table_left: float,
    table_top: float,
    cell_dim: dict,
) -> None:
    """Insert images into the table cells.

    Args:
        slide: Slide to add images to.
        table: The table containing the image cells.
        displacement_levels: List of displacement level labels.
        section_ids: List of section identifiers.
        image_dir: Directory containing experiment images.
        filename_template: Template string for image filenames.
        image_fill_uniform: If True, fit images uniformly maintaining
            aspect ratio; if False, fill available space.
        table_left: Left position of the table in inches.
        table_top: Top position of the table in inches.
        cell_dim: Dict of computed cell dimensions from
            ``_compute_cell_dimensions()``.
    """
    square_size = cell_dim["square_size"]
    available_width = cell_dim["available_width"]
    available_height = cell_dim["available_height"]

    for i, displacement in enumerate(displacement_levels):
        for j, section_id in enumerate(section_ids):
            filename = filename_template.format(
                displacement=displacement,
                section=section_id,
            )
            img_path = image_dir / filename
            for suffix in (".png", ".jpg"):
                candidate = img_path.with_suffix(suffix)
                if candidate.exists():
                    img_path = candidate
                    break
            else:
                _logger.warning(
                    "跳过不存在的图片: %s.[png|jpg]",
                    img_path,
                )
                continue

            try:
                with Image.open(img_path) as im:
                    original_width, original_height = im.size
            except Exception:
                _logger.error("无法打开图片: %s", img_path)
                continue

            aspect = original_width / original_height

            if image_fill_uniform:
                if aspect >= 1:
                    display_width = square_size
                    display_height = square_size / aspect
                else:
                    display_height = square_size
                    display_width = square_size * aspect
            else:
                if available_width / available_height >= aspect:
                    display_height = available_height
                    display_width = available_height * aspect
                else:
                    display_width = available_width
                    display_height = available_width / aspect

            image_row = 1 + i * 3
            col_a = 1 + j * 2

            cell_left = table_left + column_left_inch(table, col_a)
            cell_top = table_top + row_top_inch(table, image_row)
            cell_w = (table.columns[col_a].width + table.columns[col_a + 1].width) / Inches(1)
            cell_h = table.rows[image_row].height / Inches(1)

            offset_x = (cell_w - display_width) / 2
            offset_y = (cell_h - display_height) / 2

            left = Inches(cell_left + offset_x)
            top = Inches(cell_top + offset_y)

            try:
                slide.shapes.add_picture(
                    str(img_path),
                    left=left,
                    top=top,
                    width=Inches(display_width),
                )
            except Exception:
                _logger.error("无法插入图片: %s", img_path)
