from pptx.shapes.graphfrm import GraphicFrame
from pptx.chart.chart import Chart
from pptx.table import Table, _Cell
from pptx.slide import Slide
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pathlib import Path
from pptx import Presentation, presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml.html import etree
from PIL import Image


def main():
    IMAGE_DIR = r"testdata/images"
    OUTPUT_PPT = r"testdata/image_matrix.pptx"

    PARAM_A = ["low", "mid", "high"]
    PARAM_B = ["exp1", "exp2", "exp3", "exp4", "exp6", "exp7"]

    SUPPLEMENT = {
        ("low", "exp1"): (0.12, 0.01),
        ("low", "exp2"): (0.34, 0.02),
        ("low", "exp3"): (0.14, 0.03),
        ("low", "exp4"): (0.33, 0.04),
        ("mid", "exp1"): (0.56, 0.03),
        ("high", "exp4"): (0.78, 0.01),
    }

    prs = Presentation()

    add_summary_table_slide(
        prs,
        param_a_values=PARAM_A,
        param_b_values=PARAM_B,
        supplement_data=SUPPLEMENT,
    )

    add_image_table_slide(
        prs,
        image_dir=IMAGE_DIR,
        param_a_values=PARAM_A,
        param_b_values=PARAM_B,
        filename_template="{a}_{b}.png",
        top_left_label="卧槽牛逼",
        supplement_data=SUPPLEMENT,
        supp_row_ratio=0.30,
    )

    prs.save(OUTPUT_PPT)
    print(f"✅ PPT 已保存: {OUTPUT_PPT}")


def _convert_supplement_to_str(
    data: dict[tuple[str, str], tuple[int | float, int | float]],
) -> dict[tuple[str, str], tuple[str, str]]:
    """将数值型补充数据转换为字符串格式（第2项前加 ±）"""
    return {k: (str(v[0]), f"{v[1]}") for k, v in data.items()}


def _get_supplement_str(
    supplement_data: dict[tuple[str, str], tuple[float, float]] | None,
    a_val: str,
    b_val: str,
) -> tuple[str, str]:
    """有值时返回对应数字的字符串，缺失返回 ("", "")"""
    if supplement_data is None:
        return ("", "")
    pair = supplement_data.get((a_val, b_val))
    if pair is None:
        return ("", "")
    return (str(pair[0]), str(pair[1]))


def _create_styled_table(
    slide: Slide,
    rows: int,
    cols: int,
    left: Inches,
    top: Inches,
    width: Inches,
    height: Inches,
) -> Table:
    """创建表格并应用统一样式（黑色细框线、无填充），返回 Table 对象"""
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    _apply_table_style(table)
    return table


def _set_merged_cell(
    table: Table,
    r1: int,
    c1: int,
    r2: int,
    c2: int,
    text: str,
    bold: bool = False,
    font_size: float = 10,
    alignment: PP_ALIGN = PP_ALIGN.CENTER,
):
    """合并单元格区域并设置文本"""
    table.cell(r1, c1).merge(table.cell(r2, c2))
    _set_cell_text(
        table.cell(r1, c1), text, bold=bold, font_size=font_size, alignment=alignment
    )


def _apply_table_style(table: Table):
    """去掉表格样式和所有填充，设置黑色细框线"""
    tbl = table._tbl
    tblPr = tbl.tblPr
    if tblPr is None:
        tblPr = etree.SubElement(tbl, qn("a:tblPr"))

    # 1) 干掉表格样式 ID（这是蓝白碗的根源）
    tableStyleId = tblPr.find(qn("a:tableStyleId"))
    if tableStyleId is not None:
        tblPr.remove(tableStyleId)

    # 2) 清除表格级别的所有填充
    fill_tags = ["a:solidFill", "a:gradFill", "a:pattFill", "a:noFill", "a:grpFill"]
    for fill_tag in fill_tags:
        for el in tblPr.findall(qn(fill_tag)):
            tblPr.remove(el)

    # 3) 每个单元格：清填充 + 设细黑框线
    line_width = Pt(1)

    ln_tags = ["a:lnB", "a:lnT", "a:lnL", "a:lnR"]
    for cell in table.iter_cells():
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()

        # 清除所有类型填充
        for fill_tag in fill_tags:
            for el in tcPr.findall(qn(fill_tag)):
                tcPr.remove(el)

        # 移除旧边框
        for tag in ln_tags:
            for old in tcPr.findall(qn(tag)):
                tcPr.remove(old)

        # 添加黑色细边框
        for tag in ln_tags:
            ln = etree.SubElement(tcPr, qn(tag))
            ln.attrib["w"] = str(line_width)
            sf = etree.SubElement(ln, qn("a:solidFill"))
            srgb = etree.SubElement(sf, qn("a:srgbClr"))
            srgb.attrib["val"] = "000000"


def _set_cell_text(
    cell: _Cell,
    text: str,
    bold: bool = False,
    font_size: float = 10,
    alignment: PP_ALIGN = PP_ALIGN.CENTER,
):
    """设置表格单元格文本、字体及对齐方式"""
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = "微软雅黑"
    run.font.size = Pt(font_size)
    run.font.bold = bold


def _col_offset_inches(table: Table, col_idx: int) -> float:
    """返回从表格左边界到第 col_idx 列左边界的累计宽度（英寸）"""
    total = 0.0
    for c in range(col_idx):
        total += table.columns[c].width / Inches(1)
    return total


def _row_offset_inches(table: Table, row_idx: int) -> float:
    """返回从表格上边界到第 row_idx 行上边界的累计高度（英寸）"""
    total = 0.0
    for r in range(row_idx):
        total += table.rows[r].height / Inches(1)
    return total


def add_image_table_slide(
    prs: presentation.Presentation,
    image_dir: str,
    param_a_values: list[str],
    param_b_values: list[str],
    filename_template: str = "{a}_{b}",
    top_left_label: str = "",
    header_col_width: float = 1.2,
    header_row_height: float = 0.5,
    supplement_data: dict[tuple[str, str], tuple[float, float]] | None = None,
    supp_row_ratio: float = 0.30,
    image_fill_uniform: bool = False,
):
    image_dir: Path = Path(image_dir)
    if supplement_data is None:
        supplement_data = {}

    slide: Slide = prs.slides.add_slide(prs.slide_layouts[6])

    slide_width: float = prs.slide_width / Inches(1)
    slide_height: float = prs.slide_height / Inches(1)  # ty:ignore[unsupported-operator]

    margin_lr = 0.0
    margin_tb = 0.5
    usable_width = slide_width - 2 * margin_lr
    usable_height = slide_height - 2 * margin_tb

    n_rows = len(param_a_values)
    n_cols = len(param_b_values)
    if n_rows == 0 or n_cols == 0:
        raise ValueError("参数列表不能为空")

    data_area_width = usable_width - header_col_width
    data_area_height = usable_height - header_row_height

    total_data_cols = n_cols * 2  # 每个 b 占 2 小列
    total_data_rows = n_rows * 3  # 每个 a 占 3 小行（图片 + 标题 + 数据）

    sub_cell_width = data_area_width / total_data_cols
    sub_cell_height = data_area_height / total_data_rows

    three_row_height = 3 * sub_cell_height
    supp_total_height = three_row_height * supp_row_ratio
    img_row_height = three_row_height - supp_total_height
    title_row_height = supp_total_height / 2  # 标题行
    data_row_height = supp_total_height / 2  # 数据行

    img_area_width = sub_cell_width * 2
    img_area_height = img_row_height

    padding = Pt(1).inches
    avail_w = img_area_width - 2 * padding
    avail_h = img_area_height - 2 * padding
    square_size = min(avail_w, avail_h)

    # ---------------------------- 创建表格（使用辅助函数）----------------------------
    total_cols = 1 + total_data_cols
    total_rows = 1 + total_data_rows

    table = _create_styled_table(
        slide,
        total_rows,
        total_cols,
        Inches(margin_lr),
        Inches(margin_tb),
        Inches(usable_width),
        Inches(usable_height),
    )

    # 表格起点的英寸值（shape 的 left / top 就是 margin_lr / margin_tb）
    table_left = margin_lr
    table_top = margin_tb

    # 列宽
    table.columns[0].width = Inches(header_col_width)
    for c in range(1, total_cols):
        table.columns[c].width = Inches(sub_cell_width)

    # 行高
    table.rows[0].height = Inches(header_row_height)
    for i in range(n_rows):
        base = 1 + i * 3
        table.rows[base].height = Inches(img_row_height)  # 图片行
        table.rows[base + 1].height = Inches(title_row_height)  # 标题行
        table.rows[base + 2].height = Inches(data_row_height)  # 数据行

    # ---------------------------- 标题（使用 _set_merged_cell）----------------------------
    _set_cell_text(table.cell(0, 0), top_left_label, bold=True, font_size=12)

    for j in range(n_cols):
        col_a = 1 + j * 2
        col_b = col_a + 1
        _set_merged_cell(
            table, 0, col_a, 0, col_b, param_b_values[j], bold=True, font_size=12
        )

    for i in range(n_rows):
        row_a = 1 + i * 3
        row_b = row_a + 2
        _set_merged_cell(
            table, row_a, 0, row_b, 0, param_a_values[i], bold=True, font_size=12
        )

    # ---------------------------- 合并单元格 & 填写内容 ----------------------------
    for i, a_val in enumerate(param_a_values):
        for j, b_val in enumerate(param_b_values):
            img_row = 1 + i * 3
            title_row = img_row + 1
            data_row = img_row + 2
            col_a = 1 + j * 2
            col_b = col_a + 1

            # 图片行：两列合并（纯合并，不设文本，留给图片）
            table.cell(img_row, col_a).merge(table.cell(img_row, col_b))

            # 标题行：两列独立，填固定标题
            _set_cell_text(
                table.cell(title_row, col_a), "力值", bold=True, font_size=10
            )
            _set_cell_text(
                table.cell(title_row, col_b), "长度", bold=True, font_size=10
            )

            # 数据行：两列独立，填补充数据
            left_text, right_text = _get_supplement_str(supplement_data, a_val, b_val)
            _set_cell_text(table.cell(data_row, col_a), left_text, font_size=10)
            _set_cell_text(table.cell(data_row, col_b), right_text, font_size=10)

    # ---------------------------- 放置图片（利用表格列宽/行高定位）----------------------------
    for i, a_val in enumerate(param_a_values):
        for j, b_val in enumerate(param_b_values):
            filename = filename_template.format(a=a_val, b=b_val)
            img_path = image_dir / filename
            for suffix in (".png", ".jpg"):
                candidate = img_path.with_suffix(suffix)
                if candidate.exists():
                    img_path = candidate
                    break
            else:
                print(f"⚠️ 跳过不存在的图片: {img_path}.[png|jpg]")
                continue

            with Image.open(img_path) as im:
                orig_w, orig_h = im.size
            aspect = orig_w / orig_h

            if image_fill_uniform:
                if aspect >= 1:
                    disp_w = square_size
                    disp_h = square_size / aspect
                else:
                    disp_h = square_size
                    disp_w = square_size * aspect
            else:
                # 保持长宽比，填满可用矩形区域
                if avail_w / avail_h >= aspect:
                    disp_h = avail_h  # 高度受限
                    disp_w = avail_h * aspect
                else:
                    disp_w = avail_w  # 宽度受限
                    disp_h = avail_w / aspect

            img_row = 1 + i * 3
            col_a = 1 + j * 2

            # 通过累加列宽/行高获得单元格在幻灯片上的实际坐标
            cell_left = table_left + _col_offset_inches(table, col_a)
            cell_top = table_top + _row_offset_inches(table, img_row)
            # 图片行合并了两列，宽度取两列之和
            cell_w = (
                table.columns[col_a].width + table.columns[col_a + 1].width
            ) / Inches(1)
            cell_h = table.rows[img_row].height / Inches(1)

            offset_x = (cell_w - disp_w) / 2
            offset_y = (cell_h - disp_h) / 2

            left = Inches(cell_left + offset_x)
            top = Inches(cell_top + offset_y)

            slide.shapes.add_picture(
                str(img_path),
                left=left,
                top=top,
                width=Inches(disp_w),
            )


def add_summary_table_slide(
    prs: presentation.Presentation,
    param_a_values: list[str],
    param_b_values: list[str],
    supplement_data: dict[tuple[str, str], tuple[float, float]],  # 现在接收原始数值
    header_col_width: float = 0.9,
    param_col_width: float = 1.0,
    criteria_col_width: float = 1.2,
    header_row_height: float = 0.4,
    data_row_height: float = 0.3,
) -> None:
    """
    追加一页汇总表，并可在下方添加簇状柱形图
    """
    # ---------- 数据预处理：转成字符串供表格使用 ----------
    str_supplement = {k: (str(v[0]), f"{v[1]}") for k, v in supplement_data.items()}
    supplement_data_str = str_supplement  # 后面传递给表格部分使用
    n_rows = len(param_a_values)
    n_cols = len(param_b_values)
    total_cols = 3 + n_cols
    total_rows = 2 + n_rows * 2 + 2
    slide: Slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide_width_emu = prs.slide_width
    slide_height_emu = prs.slide_height
    # ---- 表格占宽 ----
    fixed = header_col_width + param_col_width + criteria_col_width
    total_table_inches = slide_width_emu / Inches(1)
    data_col_width = max(0, (total_table_inches - fixed) / n_cols) if n_cols else 0
    table_top = Inches(0.4)
    table_height = Inches(
        data_row_height + header_row_height + (total_rows - 2) * data_row_height
    )
    table = _create_styled_table(
        slide,
        total_rows,
        total_cols,
        Inches(0),
        table_top,
        slide_width_emu,
        table_height,
    )
    # ... 表格行高列宽设置（同原代码）...
    # 列宽
    table.columns[0].width = Inches(header_col_width)
    table.columns[1].width = Inches(param_col_width)
    table.columns[2].width = Inches(criteria_col_width)
    for c in range(3, total_cols):
        table.columns[c].width = Inches(data_col_width)
    # 行高
    table.rows[0].height = Inches(data_row_height * 1.2)
    table.rows[1].height = Inches(header_row_height)
    for r in range(2, total_rows):
        table.rows[r].height = Inches(data_row_height)
    # -------- 表格内容填充（使用 supplement_data_str）--------
    _set_merged_cell(table, 0, 0, 0, total_cols - 1, "汇总表", bold=True, font_size=14)
    _set_cell_text(table.cell(1, 0), "类别", bold=True, font_size=12)
    _set_cell_text(table.cell(1, 1), "参数", bold=True, font_size=12)
    _set_cell_text(table.cell(1, 2), "评判标准", bold=True, font_size=12)
    for j, b_val in enumerate(param_b_values):
        _set_cell_text(table.cell(1, 3 + j), b_val, bold=True, font_size=12)
    # 力值区域
    force_start, force_end = 2, 1 + n_rows
    text_value = "力值"
    if n_rows > 1:
        _set_merged_cell(
            table, force_start, 0, force_end, 0, text_value, bold=True, font_size=12
        )
    else:
        _set_cell_text(table.cell(force_start, 0), text_value, bold=True, font_size=12)
    for i, a_val in enumerate(param_a_values):
        row_idx = force_start + i
        _set_cell_text(table.cell(row_idx, 1), a_val, bold=True, font_size=12)
        _set_cell_text(table.cell(row_idx, 2), "", font_size=12)
        for j, b_val in enumerate(param_b_values):
            info = supplement_data_str.get((a_val, b_val), ("", ""))
            _set_cell_text(table.cell(row_idx, 3 + j), info[0], font_size=12)
    # 长度区域
    length_start, length_end = 2 + n_rows, 1 + n_rows * 2
    text_value = "长度"
    if n_rows > 1:
        _set_merged_cell(
            table, length_start, 0, length_end, 0, text_value, bold=True, font_size=12
        )
    else:
        _set_cell_text(table.cell(length_start, 0), text_value, bold=True, font_size=12)
    for i, a_val in enumerate(param_a_values):
        row_idx = length_start + i
        _set_cell_text(table.cell(row_idx, 1), a_val, bold=True, font_size=12)
        _set_cell_text(table.cell(row_idx, 2), "", font_size=12)
        for j, b_val in enumerate(param_b_values):
            info = supplement_data_str.get((a_val, b_val), ("", ""))
            _set_cell_text(table.cell(row_idx, 3 + j), info[1], font_size=12)
    # 结论、备注（保持原有）
    conclusion_row = 2 + n_rows * 2
    _set_cell_text(table.cell(conclusion_row, 0), "结论", bold=True, font_size=12)
    _set_merged_cell(
        table, conclusion_row, 1, conclusion_row, total_cols - 1, "", font_size=12
    )
    remark_row = conclusion_row + 1
    _set_cell_text(table.cell(remark_row, 0), "备注", bold=True, font_size=12)
    _set_merged_cell(table, remark_row, 1, remark_row, total_cols - 1, "", font_size=12)

    # ===================== 图表 =====================
    # 计算表格底部在幻灯片中的位置（英寸）
    table_bottom = (table_top + table_height) / Inches(1)  # 转换为英寸数值
    # 图表区域：从表格下方开始，占剩余高度（最大约3.5英寸，避免超出）
    chart_top_inch = table_bottom + 0.0  # 图表紧贴表格下方
    slide_height_inch = slide_height_emu / Inches(1)  # ty:ignore[unsupported-operator]
    max_chart_height = slide_height_inch - chart_top_inch - 0.2
    chart_height_inch = min(max_chart_height, 3.5)  # 限制最大高度
    chart_left_inch = 1.6
    chart_width_inch = total_table_inches - 2 * chart_left_inch
    # 构建图表数据
    chart_data = CategoryChartData()
    chart_data.categories = param_a_values  # x轴：param_A 作为类别
    for b_val in param_b_values:
        values = []
        for a_val in param_a_values:
            val = supplement_data.get((a_val, b_val), (0, 0))[0]
            values.append(val)
        chart_data.add_series(b_val, values)
    # 添加图表形状到幻灯片
    chart_frame: GraphicFrame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(chart_left_inch),
        Inches(chart_top_inch),
        Inches(chart_width_inch),
        Inches(chart_height_inch),
        chart_data,  # ty:ignore[invalid-argument-type]
    )  # ty:ignore[invalid-assignment]

    chart: Chart = chart_frame.chart

    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(10)

    chart.has_title = True
    chart.chart_title.text_frame.paragraphs[0].text = "力值对比"
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(10)

    chart.category_axis.tick_labels.font.size = Pt(10)

    chart.value_axis.has_major_gridlines = False
    chart.value_axis.tick_labels.font.size = Pt(10)

    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.show_value = True
    data_labels.font.size = Pt(10)
    for series in chart.series:
        for point in series.points:
            point.data_label.font.size = Pt(10)
            point.data_label.font.name = "微软雅黑"


if __name__ == "__main__":
    main()
